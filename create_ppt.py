import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Palette ala Anthropic: Warm, akademis, profesional
bg_color = RGBColor(253, 251, 247)  # Warm off-white/beige
text_color = RGBColor(44, 40, 38)   # Dark Charcoal/Grey
accent_color = RGBColor(217, 119, 87) # Warm terracotta

def add_slide(prs, layout_idx, title, content=[], image_placeholder_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Style Title
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = 'Georgia'
        title_shape.text_frame.paragraphs[0].font.color.rgb = accent_color
        title_shape.text_frame.paragraphs[0].font.bold = True

    # Style Body Content
    if len(slide.placeholders) > 1 and content:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for i, text in enumerate(content):
            p = tf.add_paragraph()
            p.text = text
            p.font.name = 'Segoe UI' # Clean modern sans-serif
            p.font.color.rgb = text_color
            p.font.size = Pt(24) if layout_idx == 0 else Pt(18)
            
            # If not title slide, handle indentation for bullet points
            if layout_idx != 0 and text.strip().startswith('•'):
                p.level = 1
                p.font.size = Pt(16)

    # Tambahkan placeholder area kosong untuk gambar dengan border putus-putus
    if image_placeholder_text:
        left = Inches(5.0)
        top = Inches(1.8)
        width = Inches(4.5)
        height = Inches(4.8)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = image_placeholder_text
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.font.name = 'Segoe UI'
        p.font.color.rgb = RGBColor(120, 110, 100)
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        
        # Draw dashed border
        line = txBox.line
        line.color.rgb = RGBColor(200, 190, 180)
        line.dash_style = 2 # dashed
        
        # If the slide has image, squish the text box a bit to the left so they don't overlap
        if len(slide.placeholders) > 1 and content:
            body_shape.width = Inches(4.5)

    return slide


# --- 1. TITLE SLIDE ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = bg_color

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Analisis Pola Penyebaran Narasi Politik pada Platform X"
title.text_frame.paragraphs[0].font.name = 'Georgia'
title.text_frame.paragraphs[0].font.color.rgb = accent_color
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "Pendekatan Teori Graf terhadap Data Pilpres 2024\n\nAnalisis Struktural Koordinasi Narasi (Non-Atribusi)"
for p in subtitle.text_frame.paragraphs:
    p.font.name = 'Segoe UI'
    p.font.color.rgb = text_color

# --- 2. PENDAHULUAN ---
content = [
    "Fokus Analisis:",
    "  • Membaca pola struktural jaringan (graf) berbasis kemiripan narasi dan kedekatan waktu unggah.",
    "  • Tidak bertujuan memberikan label identitas (bot, buzzer, dsb) melainkan mendeteksi anomali perilaku matematis.",
    "",
    "Konteks Metodologis:",
    "  • Memisahkan perilaku amplifikasi tunggal (Shared-Retweet) dari koordinasi penulisan narasi (Semantic Coordination).",
    "  • Retweet mudah dikenali; namun koordinasi semantik non-RT menuntut deteksi khusus (NLP)."
]
add_slide(prs, 1, "Pendahuluan & Ruang Lingkup", content)

# --- 3. METODOLOGI ---
content = [
    "Korpus Data:",
    "  • 1.002 Tweet, 584 Akun Unik, Periode Waktu 21 Menit (4 Jan 2024).",
    "",
    "Pengukuran Kemiripan Semantik (NLP):",
    "  • Menggunakan model embedding Multilingual (XLM-RoBERTa).",
    "  • Membandingkan jarak cosine antar-teks (Teks identik = Cosine 1.0).",
    "",
    "Pemodelan Temporal & Graf:",
    "  • Selisih detik dinilai dengan fungsi peluruhan eksponensial.",
    "  • Pembentukan kluster dengan Algoritma Louvain (Community Detection)."
]
add_slide(prs, 1, "Metodologi Ekstraksi Graf", content)

# --- 4. KOMPOSISI DATA (WITH IMAGE) ---
content = [
    "Proporsi Retweet Skala Masif:",
    "  • Total 56,2% dataset didominasi oleh mekanisme Retweet.",
    "",
    "Perbedaan Karakter Pola Antarkandidat:",
    "  • Anies & Prabowo dominan bergerak secara Share-Retweet.",
    "  • Ganjar minim retweet, namun volume koneksi Semantik Non-RT sangat padat (Rasio ~60:1).",
    "  • Hal ini menandakan mekanisme unggahan mandiri namun dengan narasi seragam."
]
img_placeholder = "TEMPLATE FOTO 1\n\nSilakan masukkan foto:\nTabel Komposisi Query (Tabel rasio jumlah tweet & RT untuk Anies, Ganjar, Prabowo)\n\nSumber: Tab 'Data' -> 'Komposisi Query' di visualisasi_finale.html"
add_slide(prs, 1, "Anatomi Perbincangan Politik", content, image_placeholder_text=img_placeholder)

# --- 5. KOORDINASI SEMANTIK (WITH IMAGE) ---
content = [
    "Penerapan Filter Jaringan (Backbone):",
    "  • Menyaring relasi dengan kemiripan > 0.75 dan jeda waktu sangat sempit.",
    "",
    "Temuan Utama Kluster Semantik:",
    "  • Mayoritas (kluster besar) mengkoordinasikan isu Ganjar-Mahfud (cth: nelayan, puskesmas).",
    "  • Indikasi Tinggi: Beberapa kluster memposting teks dengan similarity sempurna (1.0) dengan selisih waktu rata-rata 1 sampai 6 detik secara beruntun."
]
img_placeholder = "TEMPLATE FOTO 2\n\nSilakan masukkan foto:\nTampilan visual graf (Jaring Node warna-warni)\n\nSumber: Tab 'Ikhtisar' di visualisasi_finale.html"
add_slide(prs, 1, "Analisis Graf: Koordinasi Narasi", content, image_placeholder_text=img_placeholder)

# --- 6. PROFIL AKUN (WITH IMAGE) ---
content = [
    "Karakteristik Kuantitatif Aktor Sentral:",
    "  • Distribusi Pengikut (Followers): ~60% akun di dataset memiliki kurang dari 50 followers, dan query Ganjar memiliki median 0 follower.",
    "  • Rasio Status vs Audiens: Akun dalam kluster kuat memiliki jumlah unggahan masa lalu (Statuses) ribuan, tapi audiens nihil.",
    "",
    "Tahun Pembuatan:",
    "  • 36% dari keseluruhan akun dalam interaksi dibuat di periode 2022-2023 menjelang Pilpres."
]
img_placeholder = "TEMPLATE FOTO 3\n\nSilakan masukkan foto:\nTabel/Breakdown jumlah Followers atau bar ASCII Distribusi Followers\n\nSumber: Bagian 'Profil Akun' di Tab 'Data' visualisasi_finale.html"
add_slide(prs, 1, "Profil Aktor dalam Ekosistem", content, image_placeholder_text=img_placeholder)

# --- 7. ANOMALI TEMPORAL (WITH IMAGE) ---
content = [
    "Deteksi Spike Waktu Unggahan:",
    "  • Lonjakan ekstrem terjadi pada satu menit spesifik: 18:03 UTC.",
    "  • Sebanyak 221 tweet (22,1% dari total) dikerahkan dalam satu menit tersebut.",
    "",
    "Interpretasi Perilaku:",
    "  • Volume yang melonjak 2,3 kali lipat dari menit sebelumnya mengindikasikan kuat adanya mekanisme terjadwal (automated API posting) atau pengerahan waktu terpusat."
]
img_placeholder = "TEMPLATE FOTO 4\n\nSilakan masukkan foto:\nGrafik Bar ASCII 'Distribusi Tweet Per Menit' (yang ada highlight kuning di 18:03)\n\nSumber: Bagian 'Profil Akun' di Tab 'Data' visualisasi_finale.html"
add_slide(prs, 1, "Sinkronisasi Waktu Ekstrem", content, image_placeholder_text=img_placeholder)

# --- 8. KESIMPULAN ---
content = [
    "Pola Opini Rekayasa Terbukti:",
    "  • Keberadaan jaringan yang bukan organik berhasil diekstraksi melalui deteksi anomali pada metrik semantik (bahasa) dan temporal (waktu).",
    "",
    "Gaya Amplifikasi Berbeda:",
    "  • Kubu kandidat mengeksekusi strategi jaringan berbeda; terdapat koordinasi retweet yang jelas, namun juga praktik penulisan teks seragam yang menghindari tombol RT.",
    "",
    "Keterbatasan dan Signifikansi:",
    "  • Analisis berbasis 21 menit percakapan—menyoroti betapa padatnya strategi koordinasi di Twitter. Berguna sebagai kerangka mitigasi hoaks dan pengenalan polarisasi buatan."
]
add_slide(prs, 1, "Kesimpulan & Implikasi", content)

prs.save('Presentasi_Analisis_Narasi.pptx')
print("Presentasi berhasil dibuat: Presentasi_Analisis_Narasi.pptx")
